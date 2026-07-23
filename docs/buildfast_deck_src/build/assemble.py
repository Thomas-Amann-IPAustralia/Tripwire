from pptx import Presentation
from pptx.util import Emu
from PIL import Image
import pathlib
REND = pathlib.Path("render"); JPG = REND/"jpg"; JPG.mkdir(exist_ok=True)
SRC  = "../260610_BUILDFAST_InnoMonth.pptx"   # delivered deck, for notes
OUT  = "../260610_BUILDFAST_InnoMonth.pptx"

try:
    orig = Presentation(SRC)
    notes = [ (s.notes_slide.notes_text_frame.text if s.has_notes_slide else "") for s in orig.slides ]
except Exception:
    notes = [""]*45

prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]

for i in range(1, 46):
    png = REND/f"{i:02d}.png"
    jpg = JPG/f"{i:02d}.jpg"
    Image.open(png).convert("RGB").save(jpg, quality=93, optimize=True)
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(jpg), 0, 0, width=prs.slide_width, height=prs.slide_height)
    if i-1 < len(notes) and notes[i-1].strip():
        slide.notes_slide.notes_text_frame.text = notes[i-1]

prs.save(OUT)
print("saved", OUT)
